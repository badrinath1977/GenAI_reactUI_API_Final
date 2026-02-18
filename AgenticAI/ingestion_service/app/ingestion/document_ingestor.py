import pandas as pd
from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    TextLoader
)
from langchain.schema import Document
from pptx import Presentation


class DocumentIngestor:

    def ingest(self, file_path: str):

        ext = file_path.split('.')[-1].lower()

        # ===============================
        # PDF
        # ===============================
        if ext == "pdf":
            loader = PyPDFLoader(file_path)
            return loader.load()

        # ===============================
        # DOCX
        # ===============================
        elif ext in ["doc", "docx"]:
            loader = Docx2txtLoader(file_path)
            return loader.load()

        # ===============================
        # EXCEL
        # ===============================
        elif ext in ["xls", "xlsx"]:
            df = pd.read_excel(file_path)
            text = df.astype(str).to_string(index=False)
            return [Document(page_content=text)]

        # ===============================
        # PPT
        # ===============================
        elif ext in ["ppt", "pptx"]:
            prs = Presentation(file_path)
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            return [Document(page_content="\n".join(full_text))]

        # ===============================
        # TXT
        # ===============================
        elif ext == "txt":
            loader = TextLoader(file_path)
            return loader.load()

        else:
            raise ValueError(f"Unsupported file type: {ext}")
